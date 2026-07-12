#!/usr/bin/env python3
"""Generate minimalist Bauhaus-themed TrueComplaint AI presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Bauhaus palette (matches tailwind.config.js)
BG = RGBColor(0xF0, 0xF0, 0xF0)
FG = RGBColor(0x12, 0x12, 0x12)
RED = RGBColor(0xD0, 0x20, 0x20)
BLUE = RGBColor(0x10, 0x40, 0xC0)
YELLOW = RGBColor(0xF0, 0xC0, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xE0, 0xE0, 0xE0)

OUT = Path(__file__).resolve().parent.parent / "TrueComplaint-AI-Presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_border_box(slide, left, top, width, height, fill_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = FG
    shape.line.width = Pt(3)
    return shape


def add_text(slide, left, top, width, height, text, size=28, bold=True, color=FG, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.alignment = align
    return box


def add_label(slide, left, top, text, color=RED):
    return add_text(slide, left, top, Inches(10), Inches(0.4), text.upper(), size=11, bold=True, color=color)


def add_geometric_accent(slide):
    add_bar(slide, Inches(11.5), Inches(0.4), Inches(1.2), Inches(1.2), RED)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(5.8), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = YELLOW
    circle.line.color.rgb = FG
    circle.line.width = Pt(2)
    tri = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.2), Inches(0.55), Inches(0.55))
    tri.fill.solid()
    tri.fill.fore_color.rgb = BLUE
    tri.line.fill.background()


def slide_header(slide, label, title, dark=False):
    set_slide_bg(slide, FG if dark else BG)
    text_color = WHITE if dark else FG
    label_color = YELLOW if dark else RED
    add_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), RED if not dark else YELLOW)
    add_label(slide, Inches(0.7), Inches(0.45), label, color=label_color)
    add_text(slide, Inches(0.7), Inches(0.85), Inches(10), Inches(1.2), title.upper(), size=40, color=text_color)
    if not dark:
        add_geometric_accent(slide)


def add_bullets(slide, left, top, width, height, items, size=20, color=FG):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.font.bold = False
        p.space_after = Pt(14)
        p.level = 0


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1 — Title
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, FG)
    add_bar(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, RED)
    add_bar(s, Inches(0.35), Inches(0), Inches(0.2), SLIDE_H, BLUE)
    add_bar(s, Inches(0.55), Inches(0), Inches(0.15), SLIDE_H, YELLOW)
    add_label(s, Inches(1.2), Inches(1.8), "Ophthalmology · Voice · AI", color=YELLOW)
    add_text(s, Inches(1.2), Inches(2.3), Inches(9), Inches(1.5), "TRUECOMPLAINT AI", size=54, color=WHITE)
    add_text(s, Inches(1.2), Inches(3.5), Inches(8), Inches(0.8), "Voice Meets Medicine", size=28, bold=False, color=MUTED)
    add_text(s, Inches(1.2), Inches(4.5), Inches(9), Inches(1), "Minimal clinical intake · Structured doctor handoff", size=18, bold=False, color=WHITE)
    add_border_box(s, Inches(9.5), Inches(2), Inches(2.8), Inches(2.8), BLUE)
    add_border_box(s, Inches(10.2), Inches(2.7), Inches(1.4), Inches(1.4), WHITE)

    # 2 — Problem
    s = prs.slides.add_slide(blank)
    slide_header(s, "The Problem", "Intake Is Broken")
    add_bullets(s, Inches(0.9), Inches(2.2), Inches(5.5), Inches(4), [
        "Long paper forms patients don't understand",
        "Doctors miss critical history under time pressure",
        "Language barriers in eye clinics (Hindi / Tamil / English)",
        "No structured handoff from patient to physician",
    ], size=22)
    add_border_box(s, Inches(7.2), Inches(2.2), Inches(5), Inches(4.2), WHITE)
    add_text(s, Inches(7.5), Inches(2.5), Inches(4.4), Inches(0.5), "RESULT", size=12, color=RED)
    add_text(s, Inches(7.5), Inches(3), Inches(4.4), Inches(2.5), "Delayed care.\nLost details.\nRepeat questions.", size=26, bold=True)

    # 3 — Solution
    s = prs.slides.add_slide(blank)
    slide_header(s, "Our Solution", "Rakshak AI Intake")
    cards = [
        ("VOICE", "2–3 min call with Rakshak", RED),
        ("STRUCTURE", "Ophthalmology sheet for doctor", BLUE),
        ("EMR", "5 patients · AI assistant", YELLOW),
    ]
    for i, (title, sub, col) in enumerate(cards):
        x = Inches(0.8 + i * 4.1)
        add_border_box(s, x, Inches(2.3), Inches(3.6), Inches(3.8), WHITE)
        add_bar(s, x, Inches(2.3), Inches(3.6), Inches(0.5), col)
        add_text(s, x + Inches(0.25), Inches(3.1), Inches(3.1), Inches(0.6), title, size=22, color=FG)
        add_text(s, x + Inches(0.25), Inches(3.8), Inches(3.1), Inches(1.5), sub, size=16, bold=False)

    # 4 — Patient flow
    s = prs.slides.add_slide(blank)
    slide_header(s, "Patient Journey", "Call With Rakshak")
    steps = ["Choose language", "Live camera + mic", "7 eye intake questions", "Report sent to doctor", "Dawai voice guide"]
    for i, step in enumerate(steps):
        y = Inches(2.1 + i * 0.95)
        add_border_box(s, Inches(0.9), y, Inches(0.55), Inches(0.55), RED if i % 2 == 0 else BLUE)
        add_text(s, Inches(0.9), y + Inches(0.05), Inches(0.55), Inches(0.45), str(i + 1), size=18, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.65), y + Inches(0.05), Inches(6), Inches(0.5), step.upper(), size=20)
    add_border_box(s, Inches(8), Inches(2.2), Inches(4.5), Inches(4.5), BLUE)
    add_text(s, Inches(8.3), Inches(2.6), Inches(4), Inches(3.5), "Hindi · English · Tamil\n\nOpenAI Realtime\nWebRTC voice\n\nNatural TTS\ngpt-4o-mini-tts", size=17, bold=False, color=WHITE)

    # 5 — Ophthalmology intake
    s = prs.slides.add_slide(blank)
    slide_header(s, "Eye Clinic Intake", "Decreased Vision Protocol")
    qs = [
        "Chief complaint & duration",
        "Pain · Redness · Discharge · Floaters",
        "Onset — sudden / trauma",
        "Prior ophthalmic care",
        "Past ocular surgery",
        "Systemic disease (DM, HTN…)",
    ]
    add_bullets(s, Inches(0.9), Inches(2.1), Inches(6.5), Inches(4.5), qs, size=20)
    add_border_box(s, Inches(7.8), Inches(2.1), Inches(4.7), Inches(2), YELLOW)
    add_text(s, Inches(8.1), Inches(2.4), Inches(4.2), Inches(1.5), "DO NOT MISS\nHighlights auto-generated for physician", size=18, bold=True)

    # 6 — Doctor dashboard
    s = prs.slides.add_slide(blank)
    slide_header(s, "Physician View", "Clinical Sheet", dark=True)
    boxes = ["CC / HPI", "Associated Sx", "POH", "DDx", "Imaging · Labs"]
    for i, b in enumerate(boxes):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.1)
        y = Inches(2.2 + row * 2.2)
        add_border_box(s, x, y, Inches(3.7), Inches(1.8), WHITE)
        add_text(s, x + Inches(0.2), y + Inches(0.5), Inches(3.3), Inches(0.8), b, size=18, color=FG)
    add_text(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5), "Medical English · Urgency triage · Voice transcript", size=14, bold=False, color=MUTED)

    # 7 — EMR
    s = prs.slides.add_slide(blank)
    slide_header(s, "EMR System", "5 Hardcoded Patients")
    add_bullets(s, Inches(0.9), Inches(2.2), Inches(5.8), Inches(4), [
        "Cataract · Retinal detachment rule-out",
        "Conjunctivitis · Glaucoma F/U · Trauma",
        "AI clinical assistant (voice + text)",
        "Pre-consultation summary in seconds",
    ], size=20)
    add_border_box(s, Inches(7), Inches(2.2), Inches(5.5), Inches(4.3), WHITE)
    add_text(s, Inches(7.3), Inches(2.5), Inches(5), Inches(0.4), "SAVES TIME", size=11, color=BLUE)
    add_text(s, Inches(7.3), Inches(3), Inches(5), Inches(3), "Doctor connects with patient history before entering the room.", size=22, bold=True)

    # 8 — Dawai report
    s = prs.slides.add_slide(blank)
    slide_header(s, "Patient Column", "Dawai Voice Report")
    add_text(s, Inches(0.9), Inches(2.2), Inches(11), Inches(0.8), "Structured prescription in simple Hinglish", size=24, bold=True)
    add_bullets(s, Inches(0.9), Inches(3.1), Inches(6), Inches(3.5), [
        '"Aap ye drop din mein 3 baar daalein"',
        "Suno · Sab suno — natural OpenAI voice",
        "Precautions + follow-up read aloud",
    ], size=20)
    add_border_box(s, Inches(7.5), Inches(2.5), Inches(4.8), Inches(3.5), RED)
    add_text(s, Inches(7.8), Inches(3), Inches(4.2), Inches(2.5), "Patient\nunderstands\nmedicine", size=32, color=WHITE, align=PP_ALIGN.LEFT)

    # 9 — Tech
    s = prs.slides.add_slide(blank)
    slide_header(s, "Stack", "Architecture")
    tech = [
        ("FRONTEND", "React · Vite · Tailwind Bauhaus UI", RED),
        ("BACKEND", "Node · Express · OpenAI SDK", BLUE),
        ("VOICE", "Realtime WebRTC · Whisper · TTS cedar", YELLOW),
        ("DEPLOY", "Vercel demo UI · Localhost live AI", FG),
    ]
    for i, (t, d, c) in enumerate(tech):
        y = Inches(2 + i * 1.15)
        add_bar(s, Inches(0.9), y, Inches(0.15), Inches(0.85), c)
        add_text(s, Inches(1.2), y, Inches(2.5), Inches(0.4), t, size=14, color=c)
        add_text(s, Inches(1.2), y + Inches(0.38), Inches(10), Inches(0.5), d, size=18, bold=False)

    # 10 — Close
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BLUE)
    add_bar(s, Inches(0), Inches(6.9), SLIDE_W, Inches(0.12), YELLOW)
    add_text(s, Inches(0), Inches(2.5), SLIDE_W, Inches(1.2), "THANK YOU", size=52, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(3.8), SLIDE_W, Inches(0.8), "TrueComplaint AI · Rakshak · Eye Clinic Intake", size=20, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(5), SLIDE_W, Inches(0.6), "Demo: localhost  ·  UI: Vercel", size=14, bold=False, color=MUTED, align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
